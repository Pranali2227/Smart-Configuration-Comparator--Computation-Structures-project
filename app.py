from flask import Flask, request, send_from_directory, jsonify
import os
import tempfile
import difflib
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
from PIL import Image
import io

# Configure pytesseract path if needed (for Windows)
pytesseract.pytesseract.tesseract_cmd = r'C:\Users\patil\OneDrive\Desktop\College DA\Computation Structures\Project\tesseract.exe'  # Adjust path as necessary

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'static/uploads'

def extract_text_from_file(file_path, file_ext):
    text = ""
    if file_ext == '.pdf':
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:  # Only add if text is extracted
                    text += page_text + "\n"
    elif file_ext in ['.doc', '.docx']:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_ext in ['.xls', '.xlsx']:
        wb = openpyxl.load_workbook(file_path)
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell])
                if row_text.strip():
                    text += row_text + "\n"
    elif file_ext in ['.ppt', '.pptx']:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file_ext in ['.txt', '.rtf', '.html']:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
    elif file_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img, config='--psm 6 --oem 3')
        except pytesseract.pytesseract.TesseractNotFoundError:
            text = "Tesseract OCR not installed. Please install Tesseract from https://github.com/UB-Mannheim/tesseract/wiki and ensure it's in your PATH."
        except Exception as e:
            text = f"Error extracting text from image: {str(e)}"
    return text.strip()

def format_tabular(text):
    lines = text.split('\n')
    formatted = []
    for line in lines:
        if line.strip():
            parts = line.split()
            if len(parts) > 1:
                formatted.append('| ' + ' | '.join(parts) + ' |')
            else:
                formatted.append(line)
    return '\n'.join(formatted)

def compare_texts(text1, text2):
    # Normalize text: convert to lowercase for case-insensitive comparison
    normalize = lambda text: text.lower().strip()
    norm_text1 = normalize(text1)
    norm_text2 = normalize(text2)

    # Use difflib for better line-based diff
    diff = list(difflib.unified_diff(norm_text1.splitlines(keepends=True), norm_text2.splitlines(keepends=True), fromfile='File 1', tofile='File 2', lineterm=''))
    similarities = []
    file1_diff = []
    file2_diff = []
    for line in diff:
        if line.startswith(' '):
            similarities.append(line[1:])
        elif line.startswith('-'):
            file1_diff.append(line)
        elif line.startswith('+'):
            file2_diff.append(line)
    similarities = '\n'.join(similarities).strip()
    file1_diff = '\n'.join(file1_diff).strip()
    file2_diff = '\n'.join(file2_diff).strip()
    return similarities, file1_diff, file2_diff

@app.route('/')
def index():
    return send_from_directory('.', 'index.html')

@app.route('/<path:filename>')
def serve_static(filename):
    if filename.startswith('static/'):
        return send_from_directory('static', filename[7:])
    return send_from_directory('.', filename)

@app.route('/compare', methods=['POST'])
def compare():
    file1 = request.files['file1']
    file2 = request.files['file2']
    if not file1 or not file2:
        return jsonify({'error': 'Both files are required'}), 400

    ext1 = os.path.splitext(file1.filename)[1].lower()
    ext2 = os.path.splitext(file2.filename)[1].lower()

    file1_path = os.path.join(app.config['UPLOAD_FOLDER'], file1.filename)
    file2_path = os.path.join(app.config['UPLOAD_FOLDER'], file2.filename)
    file1.save(file1_path)
    file2.save(file2_path)

    try:
        text1 = extract_text_from_file(file1_path, ext1)
        text2 = extract_text_from_file(file2_path, ext2)
        similarities, file1_diff, file2_diff = compare_texts(text1, text2)
        file1_url = f"/static/uploads/{file1.filename}" if ext1 in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff'] else None
        file2_url = f"/static/uploads/{file2.filename}" if ext2 in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff'] else None
        # Format tabular data for XLSX
        if ext1 in ['.xls', '.xlsx']:
            text1 = format_tabular(text1)
        if ext2 in ['.xls', '.xlsx']:
            text2 = format_tabular(text2)

        return jsonify({
            'file1_content': text1,
            'file2_content': text2,
            'file1_url': file1_url,
            'file2_url': file2_url,
            'similarities': similarities,
            'file1_diff': file1_diff,
            'file2_diff': file2_diff
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        # Do not remove files for images to keep them accessible
        if not ext1 in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            if os.path.exists(file1_path):
                os.remove(file1_path)
        if not ext2 in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            if os.path.exists(file2_path):
                os.remove(file2_path)

if __name__ == '__main__':
    app.run(debug=True)
